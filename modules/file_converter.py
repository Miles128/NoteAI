import re
import shutil
import subprocess
import tempfile
from abc import ABC, abstractmethod
from collections.abc import Callable
from pathlib import Path

from config.settings import RAW_FOLDER
from utils.helpers import (
    clean_text,
    ensure_dir,
    get_file_extension,
    read_file_with_encoding,
    remove_images_from_markdown,
    sanitize_filename,
)
from utils.logger import logger
from utils.pdf_utils import extract_pdf_pages, extract_pdf_text
from utils.tag_extractor import (
    add_yaml_frontmatter_to_content,
    add_yaml_frontmatter_to_file,
    extract_tags_from_filename,
)


class BaseConverter(ABC):
    """转换器基类"""

    def __init__(self, progress_callback: Callable | None = None):
        self.progress_callback = progress_callback

    @abstractmethod
    def to_markdown(self, file_path: str) -> str:
        """转换为Markdown"""
        pass


class PDFConverter(BaseConverter):
    """PDF转Markdown转换器（仅使用快速路径）"""

    SUPPORTED_FORMATS = [".pdf"]

    MIN_PAGE_COUNT_FOR_SIGNATURE_DETECTION = 3
    SIGNATURE_MIN_LENGTH = 5
    SIGNATURE_MAX_LENGTH = 200
    SIGNATURE_MAX_LINES = 5

    def to_markdown(self, file_path: str) -> str:
        """使用 PyMuPDF 将 PDF 转换为 Markdown"""
        logger.info(f"开始转换PDF: {file_path}")

        try:
            markdown_content = self._extract_pdf_text(file_path)
            markdown_content = clean_text(markdown_content)
            markdown_content = self._remove_repeated_signatures(markdown_content, file_path)

            markdown_content = remove_images_from_markdown(markdown_content)

            logger.info(f"PDF转换完成: {file_path}")
            return markdown_content

        except Exception as e:
            logger.error(f"PDF转换失败 {file_path}: {e}")
            raise

    def _extract_pdf_text(self, file_path: str) -> str:
        """使用 PyMuPDF 提取 PDF 文本（统一入口）"""
        return extract_pdf_text(file_path)

    def _extract_page_texts(self, file_path: str) -> list[str]:
        """提取每一页的文本用于签名检测"""
        return extract_pdf_pages(file_path)

    def _find_signature_lines(self, page_texts: list[str]) -> set:
        """
        识别在每一页都出现的重复文本行（签名、页眉、页脚等）

        Returns:
            包含重复文本的集合
        """
        if len(page_texts) < self.MIN_PAGE_COUNT_FOR_SIGNATURE_DETECTION:
            return set()

        signature_candidates = []
        for page_text in page_texts:
            lines = page_text.split("\n")
            lines = [line.strip() for line in lines if line.strip()]
            signature_candidates.append(set(lines))

        common_lines = signature_candidates[0]
        for candidate_set in signature_candidates[1:]:
            common_lines &= candidate_set

        valid_signatures = set()
        for line in common_lines:
            if self.SIGNATURE_MIN_LENGTH <= len(line) <= self.SIGNATURE_MAX_LENGTH:
                valid_signatures.add(line)

        logger.info(f"检测到 {len(valid_signatures)} 条重复签名/页眉/页脚")
        return valid_signatures

    def _remove_repeated_signatures(self, markdown_content: str, file_path: str) -> str:
        """从Markdown内容中移除重复出现的签名、页眉、页脚（单次打开PDF）"""
        import fitz

        with fitz.open(file_path) as doc:
            page_count = len(doc)
            page_texts = [page.get_text() for page in doc]

        if page_count < self.MIN_PAGE_COUNT_FOR_SIGNATURE_DETECTION:
            logger.info(f"页数({page_count})少于{self.MIN_PAGE_COUNT_FOR_SIGNATURE_DETECTION}，跳过签名检测")
            return markdown_content

        signature_lines = self._find_signature_lines(page_texts)

        if not signature_lines:
            return markdown_content

        lines = markdown_content.split("\n")
        filtered_lines = []
        skip_count = 0

        for line in lines:
            stripped = line.strip()
            if stripped in signature_lines:
                skip_count += 1
                continue
            filtered_lines.append(line)

        if skip_count > 0:
            logger.info(f"已移除 {skip_count} 行签名/页眉/页脚内容")

        return "\n".join(filtered_lines)


class TXTConverter(BaseConverter):
    def to_markdown(self, file_path: str) -> str:
        """TXT转Markdown"""
        logger.info(f"开始转换TXT: {file_path}")

        try:
            raw_content = read_file_with_encoding(file_path)

            markdown_content = clean_text(raw_content)

            markdown_content = remove_images_from_markdown(markdown_content)

            logger.info(f"TXT转换完成: {file_path}")
            return markdown_content

        except Exception as e:
            logger.error(f"TXT转换失败: {e}")
            raise


class DOCXConverter(BaseConverter):
    """Word文档转Markdown转换器（使用mammoth，仅支持 .docx）"""

    SUPPORTED_FORMATS = [".docx"]

    def to_markdown(self, file_path: str) -> str:
        """将Word文档转换为Markdown"""
        logger.info(f"开始转换Word文档: {file_path}")

        try:
            markdown_content = self._extract_docx_text(file_path)
            markdown_content = clean_text(markdown_content)

            markdown_content = remove_images_from_markdown(markdown_content)

            logger.info(f"Word文档转换完成: {file_path}")
            return markdown_content

        except Exception as e:
            logger.error(f"Word文档转换失败 {file_path}: {e}")
            raise

    def _extract_docx_text(self, file_path: str) -> str:
        """使用mammoth将DOCX转换为Markdown"""
        import mammoth

        with open(file_path, "rb") as docx_file:
            result = mammoth.convert_to_markdown(docx_file)
            return result.value


class LegacyDOCConverter(BaseConverter):
    """旧版 Word .doc 转 Markdown，依赖系统可用的文本提取工具。"""

    SUPPORTED_FORMATS = [".doc"]

    def to_markdown(self, file_path: str) -> str:
        logger.info(f"开始转换旧版Word文档: {file_path}")

        errors = []
        for extractor in (self._extract_with_textutil, self._extract_with_antiword, self._extract_with_catdoc):
            try:
                content = extractor(file_path)
                if content and content.strip():
                    markdown_content = clean_text(content)
                    markdown_content = remove_images_from_markdown(markdown_content)
                    logger.info(f"旧版Word文档转换完成: {file_path}")
                    return markdown_content
            except Exception as e:
                errors.append(str(e))

        msg = "无法转换旧版 .doc 文件：需要 macOS textutil，或安装 antiword/catdoc"
        if errors:
            msg += f"；错误: {'; '.join(errors[-2:])}"
        raise RuntimeError(msg)

    def _extract_with_textutil(self, file_path: str) -> str:
        textutil = shutil.which("textutil")
        if not textutil:
            raise RuntimeError("textutil not found")

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / (Path(file_path).stem + ".txt")
            proc = subprocess.run(
                [textutil, "-convert", "txt", "-output", str(out_path), file_path],
                capture_output=True,
                text=True,
                timeout=60,
                check=False,
            )
            if proc.returncode != 0:
                raise RuntimeError((proc.stderr or proc.stdout or "textutil failed").strip())
            return read_file_with_encoding(str(out_path))

    def _extract_with_antiword(self, file_path: str) -> str:
        antiword = shutil.which("antiword")
        if not antiword:
            raise RuntimeError("antiword not found")
        proc = subprocess.run(
            [antiword, file_path],
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )
        if proc.returncode != 0:
            raise RuntimeError((proc.stderr or "antiword failed").strip())
        return proc.stdout

    def _extract_with_catdoc(self, file_path: str) -> str:
        catdoc = shutil.which("catdoc")
        if not catdoc:
            raise RuntimeError("catdoc not found")
        proc = subprocess.run(
            [catdoc, file_path],
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )
        if proc.returncode != 0:
            raise RuntimeError((proc.stderr or "catdoc failed").strip())
        return proc.stdout


class PPTConverter(BaseConverter):
    """PPT转Markdown转换器（使用python-pptx，仅支持 .pptx）"""

    SUPPORTED_FORMATS = [".pptx"]

    def to_markdown(self, file_path: str) -> str:
        logger.info(f"开始转换PPT: {file_path}")

        try:
            markdown_content = self._extract_pptx_text(file_path)
            markdown_content = clean_text(markdown_content)

            markdown_content = remove_images_from_markdown(markdown_content)

            logger.info(f"PPT转换完成: {file_path}")
            return markdown_content

        except Exception as e:
            logger.error(f"PPT转换失败 {file_path}: {e}")
            raise

    def _extract_pptx_text(self, file_path: str) -> str:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"## 幻灯片 {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                    if rows:
                        header = rows[0]
                        sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
                        parts.append(header)
                        parts.append(sep)
                        parts.extend(rows[1:])
            if len(parts) > 1:
                slides.append("\n\n".join(parts))
        return "\n\n---\n\n".join(slides)


class LegacyPPTConverter(BaseConverter):
    """旧版 PowerPoint .ppt 转 Markdown，解析 OLE PowerPoint Document 文本记录。"""

    SUPPORTED_FORMATS = [".ppt"]
    TEXT_CHARS_ATOM = 4000
    TEXT_BYTES_ATOM = 4008
    CSTRING = 4026

    def to_markdown(self, file_path: str) -> str:
        logger.info(f"开始转换旧版PPT: {file_path}")

        try:
            markdown_content = self._extract_ppt_text(file_path)
            markdown_content = clean_text(markdown_content)
            markdown_content = remove_images_from_markdown(markdown_content)
            logger.info(f"旧版PPT转换完成: {file_path}")
            return markdown_content
        except Exception as e:
            logger.error(f"旧版PPT转换失败 {file_path}: {e}")
            raise

    def _extract_ppt_text(self, file_path: str) -> str:
        import olefile

        with olefile.OleFileIO(file_path) as ole:
            stream_name = self._find_powerpoint_stream(ole)
            if not stream_name:
                raise RuntimeError("未找到 PowerPoint Document 数据流")
            data = ole.openstream(stream_name).read()

        texts = self._extract_text_records(data)
        if not texts:
            raise RuntimeError("未能从旧版 .ppt 文件中提取文本")
        return "\n\n".join(texts)

    def _find_powerpoint_stream(self, ole) -> list[str] | None:
        for entry in ole.listdir(streams=True, storages=False):
            if entry and entry[-1] == "PowerPoint Document":
                return entry
        return None

    def _extract_text_records(self, data: bytes) -> list[str]:
        texts = []
        pos = 0
        data_len = len(data)
        while pos + 8 <= data_len:
            rec_header = int.from_bytes(data[pos : pos + 2], "little")
            rec_type = int.from_bytes(data[pos + 2 : pos + 4], "little")
            rec_len = int.from_bytes(data[pos + 4 : pos + 8], "little")
            next_pos = pos + 8 + rec_len
            if rec_len < 0 or next_pos > data_len:
                pos += 1
                continue

            payload = data[pos + 8 : next_pos]
            if rec_type == self.TEXT_CHARS_ATOM:
                text = self._decode_utf16(payload)
                if text:
                    texts.append(text)
            elif rec_type == self.TEXT_BYTES_ATOM:
                text = self._decode_latin_text(payload)
                if text:
                    texts.append(text)
            elif rec_type == self.CSTRING:
                text = self._decode_cstring(payload)
                if text:
                    texts.append(text)

            pos = next_pos if next_pos > pos else pos + 1
            _ = rec_header

        return self._dedupe_texts(texts)

    def _decode_utf16(self, payload: bytes) -> str:
        try:
            return self._normalize_text(payload.decode("utf-16le", errors="ignore"))
        except Exception:
            return ""

    def _decode_latin_text(self, payload: bytes) -> str:
        return self._normalize_text(payload.decode("latin-1", errors="ignore"))

    def _decode_cstring(self, payload: bytes) -> str:
        utf16 = self._decode_utf16(payload)
        latin = self._decode_latin_text(payload)
        return utf16 if self._text_score(utf16) >= self._text_score(latin) else latin

    def _normalize_text(self, text: str) -> str:
        text = text.replace("\x00", "")
        lines = [line.strip() for line in re.split(r"[\r\n]+", text) if line.strip()]
        return "\n".join(lines).strip()

    def _text_score(self, text: str) -> int:
        if not text:
            return 0
        score = len(text)
        score -= sum(4 for ch in text if ord(ch) < 32 and ch not in "\n\r\t")
        score += sum(2 for ch in text if "\u4e00" <= ch <= "\u9fff")
        score += sum(1 for ch in text if ch.isalpha() or ch.isdigit())
        return score

    def _dedupe_texts(self, texts: list[str]) -> list[str]:
        deduped = []
        seen = set()
        for text in texts:
            key = re.sub(r"\s+", " ", text).strip()
            if len(key) < 2 or key in seen:
                continue
            seen.add(key)
            deduped.append(text)
        return deduped


class HTMLConverter(BaseConverter):
    """HTML转Markdown转换器（使用html2text）"""

    SUPPORTED_FORMATS = [".html", ".htm"]

    def to_markdown(self, file_path: str) -> str:
        logger.info(f"开始转换HTML: {file_path}")

        try:
            markdown_content = self._extract_html_text(file_path)
            markdown_content = clean_text(markdown_content)

            markdown_content = remove_images_from_markdown(markdown_content)

            logger.info(f"HTML转换完成: {file_path}")
            return markdown_content

        except Exception as e:
            logger.error(f"HTML转换失败 {file_path}: {e}")
            raise

    def _extract_html_text(self, file_path: str) -> str:
        import html2text

        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False
        h.body_width = 0
        html_content = read_file_with_encoding(file_path)
        return h.handle(html_content)


class FileConverterManager:
    """文件转换管理器"""

    PDF_FORMATS = [".pdf"]
    DOCX_FORMATS = [".docx"]
    LEGACY_DOC_FORMATS = [".doc"]
    PPT_FORMATS = [".pptx"]
    LEGACY_PPT_FORMATS = [".ppt"]
    HTML_FORMATS = [".html", ".htm"]
    TXT_FORMATS = [".txt"]

    def __init__(self, progress_callback: Callable | None = None):
        self.progress_callback = progress_callback
        self._pdf_converter = None
        self._docx_converter = None
        self._legacy_doc_converter = None
        self._txt_converter = None
        self._ppt_converter = None
        self._legacy_ppt_converter = None
        self._html_converter = None

    @staticmethod
    def _needs_llm_rewrite(content: str) -> bool:
        if not content or len(content.strip()) < 50:
            return True
        text = content.strip()
        has_heading = bool(re.search(r"^#{1,3}\s+\S", text, re.MULTILINE))
        has_sentence_end = bool(re.search(r"[。！？.!?\n]", text))
        has_comma = bool(re.search(r"[，,、；;：:]", text))
        has_paragraph = text.count("\n\n") >= 1
        if has_heading and has_sentence_end:
            return False
        if has_sentence_end and has_comma and len(text) > 100:
            return False
        if has_paragraph and has_sentence_end:
            return False
        return True

    @property
    def pdf_converter(self):
        """懒加载PDF转换器"""
        if self._pdf_converter is None:
            self._pdf_converter = PDFConverter(self.progress_callback)
        return self._pdf_converter

    @property
    def docx_converter(self):
        """懒加载Word文档转换器"""
        if self._docx_converter is None:
            self._docx_converter = DOCXConverter(self.progress_callback)
        return self._docx_converter

    @property
    def legacy_doc_converter(self):
        if self._legacy_doc_converter is None:
            self._legacy_doc_converter = LegacyDOCConverter(self.progress_callback)
        return self._legacy_doc_converter

    @property
    def txt_converter(self):
        """懒加载TXT转换器"""
        if self._txt_converter is None:
            self._txt_converter = TXTConverter(self.progress_callback)
        return self._txt_converter

    @property
    def ppt_converter(self):
        if self._ppt_converter is None:
            self._ppt_converter = PPTConverter(self.progress_callback)
        return self._ppt_converter

    @property
    def legacy_ppt_converter(self):
        if self._legacy_ppt_converter is None:
            self._legacy_ppt_converter = LegacyPPTConverter(self.progress_callback)
        return self._legacy_ppt_converter

    @property
    def html_converter(self):
        if self._html_converter is None:
            self._html_converter = HTMLConverter(self.progress_callback)
        return self._html_converter

    def _get_converter(self, ext: str):
        """根据扩展名获取合适的转换器"""
        ext = ext.lower()
        if ext in self.PDF_FORMATS:
            return self.pdf_converter
        if ext in self.DOCX_FORMATS:
            return self.docx_converter
        if ext in self.LEGACY_DOC_FORMATS:
            return self.legacy_doc_converter
        if ext in self.PPT_FORMATS:
            return self.ppt_converter
        if ext in self.LEGACY_PPT_FORMATS:
            return self.legacy_ppt_converter
        if ext in self.HTML_FORMATS:
            return self.html_converter
        if ext in self.TXT_FORMATS:
            return self.txt_converter
        return None

    def convert_file(
        self,
        file_path: str,
        output_path: str,
        output_format: str = "markdown",
    ) -> dict:
        """
        转换单个文件

        Args:
            file_path: 输入文件路径
            output_path: 输出目录路径
            output_format: 输出格式（目前仅支持 markdown）
        """
        result = {"file_path": file_path, "success": False, "output_path": None, "error": None}

        try:
            file_path_obj = Path(file_path)

            if not file_path_obj.exists():
                result["error"] = "文件不存在"
                return result

            ext = get_file_extension(file_path_obj.name)
            converter = self._get_converter(ext)

            if converter is None:
                result["error"] = f"不支持的文件格式: {ext}"
                return result

            # 转换为Markdown
            markdown_content = converter.to_markdown(str(file_path))

            if self._needs_llm_rewrite(markdown_content):
                try:
                    from utils.llm_utils import call_llm_raw, check_api_config

                    ok, msg = check_api_config()
                    if ok:
                        from prompts import GENERAL_CONTENT_TO_MARKDOWN_PROMPT

                        prompt = GENERAL_CONTENT_TO_MARKDOWN_PROMPT.format(
                            source_type=ext, content=markdown_content[:6000]
                        )
                        rewritten = call_llm_raw(prompt, temperature=0.3)
                        if rewritten and len(rewritten.strip()) > len(markdown_content.strip()) * 0.3:
                            markdown_content = rewritten.strip()
                            logger.info(f"LLM 重写转换结果: {file_path}")
                except Exception as e:
                    logger.warning(f"LLM 重写失败，保留原始转换: {e}")

            # 添加YAML front matter
            markdown_content = add_yaml_frontmatter_to_content(markdown_content, tags=[], source=file_path)

            # 保存文件
            output_dir = ensure_dir(output_path)
            output_filename = sanitize_filename(file_path_obj.stem) + ".md"
            output_file = output_dir / output_filename

            counter = 1
            original_output_file = output_file
            while output_file.exists():
                stem = original_output_file.stem
                output_file = original_output_file.parent / f"{stem}_{counter}.md"
                counter += 1

            with open(output_file, "w", encoding="utf-8") as f:
                f.write(markdown_content)

            tags = extract_tags_from_filename(str(output_file))
            if tags:
                add_yaml_frontmatter_to_file(str(output_file), tags=tags, source=file_path)

            try:
                from utils.topic_assigner import auto_assign_topic_for_file

                auto_assign_topic_for_file(str(output_file))
            except Exception as e:
                logger.warning(f"自动分配主题失败: {e}")

            result["success"] = True
            result["output_path"] = str(output_file)
            result["tags"] = tags

            logger.info(f"文件转换成功: {output_file}")

        except Exception as e:
            result["error"] = str(e)
            logger.error(f"文件转换失败: {e}")

        return result

    def convert_batch(
        self,
        file_paths: list[str],
        output_path: str,
        raw_path: str = None,
        output_format: str = "markdown",
    ) -> list[dict]:
        """批量转换文件"""
        results = []
        total = len(file_paths)

        for i, file_path in enumerate(file_paths):
            if self.progress_callback:
                self.progress_callback(i + 1, total, f"正在转换: {Path(file_path).name}")

            result = self.convert_file(file_path, output_path, output_format)

            if result["success"] and raw_path:
                self._move_to_raw(file_path, raw_path)

            results.append(result)

        if self.progress_callback:
            success_count = sum(1 for r in results if r["success"])
            self.progress_callback(total, total, f"转换完成: {success_count}/{total} 成功")

        return results

    def _move_to_raw(self, file_path: str, raw_path: str) -> bool:
        """移动文件到Raw文件夹进行归档"""
        try:
            source = Path(file_path)
            if not source.exists():
                logger.warning(f"源文件不存在，跳过移动: {file_path}")
                return False

            raw_dir = Path(raw_path)
            raw_dir.mkdir(parents=True, exist_ok=True)

            dest = raw_dir / source.name
            if dest.exists():
                base_name = source.stem
                extension = source.suffix
                counter = 1
                while dest.exists():
                    dest = raw_dir / f"{base_name}_{counter}{extension}"
                    counter += 1

            import shutil

            shutil.move(str(source), str(dest))
            logger.info(f"已移动文件到Raw: {source} -> {dest}")
            return True
        except Exception as e:
            logger.error(f"移动文件到Raw失败: {file_path}, 错误: {e}")
            return False

    def convert_folder(
        self,
        folder_path: str,
        output_path: str,
        raw_path: str = None,
        output_format: str = "markdown",
        recursive: bool = True,
    ) -> list[dict]:
        """转换整个文件夹"""
        folder = Path(folder_path)

        if not folder.exists():
            logger.error(f"文件夹不存在: {folder_path}")
            return [
                {
                    "file_path": folder_path,
                    "success": False,
                    "output_path": None,
                    "error": f"文件夹不存在: {folder_path}",
                }
            ]

        file_paths = []
        all_files = list(folder.rglob("*") if recursive else folder.glob("*"))

        for file_path in all_files:
            if file_path.is_file() and not file_path.name.startswith("."):
                relative_path = file_path.relative_to(folder)
                if RAW_FOLDER in relative_path.parts:
                    continue
                ext = file_path.suffix.lower()
                if ext != ".md" and ext in self.get_supported_formats():
                    file_paths.append(str(file_path))

        logger.info(f"找到 {len(file_paths)} 个可转换文件")

        return self.convert_batch(file_paths, output_path, raw_path, output_format)

    @classmethod
    def get_supported_formats(cls) -> list[str]:
        """获取支持的格式列表"""
        return (
            cls.PDF_FORMATS
            + cls.DOCX_FORMATS
            + cls.LEGACY_DOC_FORMATS
            + cls.PPT_FORMATS
            + cls.LEGACY_PPT_FORMATS
            + cls.HTML_FORMATS
            + cls.TXT_FORMATS
        )
